import streamlit as st
from PIL import Image, ImageDraw, ImageFont
import pandas as pd
from io import BytesIO
import zipfile
import re
import os

# --- นำเข้าไลบรารีสำหรับคลิกหาพิกัด ---
try:
    from streamlit_image_coordinates import streamlit_image_coordinates
except ImportError:
    st.error("⚠️ ไม่พบไลบรารี streamlit-image-coordinates")
    st.info("กรุณาเปิด Terminal แล้วพิมพ์: pip install streamlit-image-coordinates")
    st.stop()

# --- นำเข้าไลบรารีสำหรับสร้าง PowerPoint และจัดการข้อความ ---
try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    st.error("⚠️ ไม่พบไลบรารี python-pptx (สำหรับสร้าง PowerPoint)")
    st.info("กรุณาเปิด Terminal แล้วพิมพ์: pip install python-pptx")
    st.stop()

# ==========================================
# 🛠️ HELPER FUNCTIONS
# ==========================================

def get_system_font_path():
    """ค้นหาฟอนต์ที่ปรับขนาดได้ในระบบเพื่อใช้เป็นสำรอง (จากโค้ดที่ 2)"""
    paths = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
        "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
        "C:\\Windows\\Fonts\\arial.ttf",
        "C:\\Windows\\Fonts\\tahoma.ttf",
        "/System/Library/Fonts/Helvetica.ttf"
    ]
    for p in paths:
        if os.path.exists(p):
            return p
    return None

def fix_thai_text(text):
    """
    จัดตำแหน่งสระและวรรณยุกต์สำหรับ 'ฟอนต์รุ่นเก่า' (ใช้รหัส PUA)
    รวมถึงการตัดฐาน ญ และ ฐ เมื่อมีสระล่าง (ครอบคลุม fix_thai_baseless_chars จากโค้ดที่ 2)
    """
    if not isinstance(text, str):
        return str(text) if pd.notna(text) else ""
        
    tone_marks = ['\u0e48', '\u0e49', '\u0e4a', '\u0e4b', '\u0e4c']
    
    # 1. วรรณยุกต์ที่ตามหลังสระบน
    upper_vowels = ['\u0e31', '\u0e34', '\u0e35', '\u0e36', '\u0e37', '\u0e4d']
    high_tone_marks = ['\uf713', '\uf714', '\uf715', '\uf716', '\uf717']
    
    for i, tone in enumerate(tone_marks):
        for vowel in upper_vowels:
            text = text.replace(vowel + tone, vowel + high_tone_marks[i])
    
    # 2. วรรณยุกต์ที่ตามหลัง ป ฝ ฟ
    tall_consonants = ['ป', 'ฝ', 'ฟ']
    left_tone_marks = ['\uf70a', '\uf70b', '\uf70c', '\uf70d', '\uf70e']
    
    for i, tone in enumerate(tone_marks):
        for cons in tall_consonants:
            text = text.replace(cons + tone, cons + left_tone_marks[i])
            
    # 3. สระอำ
    text = text.replace('\u0e4d\u0e32', '\u0e33')
    
    # 4. ญ และ ฐ เมื่อมีสระล่าง (ตัดฐาน)
    text = text.replace('ญุ', '\uf70fุ').replace('ญู', '\uf70fู').replace('ญฺ', '\uf70fฺ')
    text = text.replace('ฐุ', '\uf700ุ').replace('ฐู', '\uf700ู').replace('ฐฺ', '\uf700ฺ')
    
    return text

def get_font(font_name, size):
    """โหลดฟอนต์จากไฟล์ที่ผู้ใช้อัปโหลด (รองรับหลายฟอนต์จากโค้ดที่ 2) หรือฟอนต์ระบบ"""
    # 1. ค้นหาฟอนต์ที่ระบุจาก Dictionary
    if font_name and 'fonts_dict' in st.session_state and font_name in st.session_state.fonts_dict:
        try:
            return ImageFont.truetype(BytesIO(st.session_state.fonts_dict[font_name]), size)
        except Exception as e:
            st.error(f"❌ ไม่สามารถโหลดฟอนต์ '{font_name}': {e}")

    # 2. ฟอนต์เดี่ยวสำรอง (จากโค้ดแรกดั้งเดิม)
    if 'font_bytes' in st.session_state and st.session_state.font_bytes:
        try:
            return ImageFont.truetype(BytesIO(st.session_state.font_bytes), size)
        except Exception as e:
            pass

    # 3. หากไม่พบ พยายามใช้ฟอนต์ระบบที่ปรับขนาดได้ (จากโค้ดที่ 2)
    sys_path = get_system_font_path()
    if sys_path:
        return ImageFont.truetype(sys_path, size)
        
    return ImageFont.load_default()

def sanitize_filename(name):
    return re.sub(r'[<>:"/\\|?*]', '_', str(name)).strip() or "certificate"

def render_certificate(template_img, texts, row_data=None, font_version="ใหม่"):
    img = template_img.copy()
    if img.mode != 'RGB':
        img = img.convert('RGB')

    draw = ImageDraw.Draw(img)
    is_new_font = font_version == "ใหม่"

    for txt in texts:
        if txt['type'] == 'static':
            content = txt['text']
        else:
            if row_data and txt['column'] in row_data:
                val = row_data[txt['column']]
                content = str(val) if pd.notna(val) else ""
            else:
                content = "ตัวอย่างข้อมูล"

        if not content: continue

        # ถ้าเป็นฟอนต์เก่า ให้ใช้โค้ดดัดแปลง PUA ถ้าเป็นฟอนต์ใหม่ไม่ต้องแปลง
        if not is_new_font:
            content = fix_thai_text(content)

        # ใช้ฟอนต์ตามชื่อที่ระบุในแต่ละข้อความ (จากโค้ดที่ 2)
        font = get_font(txt.get('font_name'), txt['size'])
        
        # === วัดขนาดและวาดข้อความ (คงโค้ดแรกไว้เพราะรองรับภาษาไทยรุ่นใหม่ได้ดีกว่า) ===
        
        # วัดขนาดความกว้าง
        if is_new_font:
            try:
                # ลองใช้ระบบจัดเรียงสระแบบใหม่
                text_width = draw.textlength(content, font=font, language="th")
            except Exception: 
                # ถ้าเซิร์ฟเวอร์ไม่รองรับ (เกิด KeyError, TypeError ฯลฯ) ให้ใช้แบบธรรมดา
                text_width = draw.textlength(content, font=font)
        else:
            text_width = draw.textlength(content, font=font)

        # หาจุด X ทางซ้ายสุดเพื่อให้ข้อความอยู่กึ่งกลางพอยน์ที่คลิก
        start_x = txt['x'] - (text_width / 2)

        # วาดข้อความ
        if is_new_font:
            try:
                draw.text((start_x, txt['y']), content, fill=txt['color'], font=font, anchor="ls", language="th")
            except Exception:
                draw.text((start_x, txt['y']), content, fill=txt['color'], font=font, anchor="ls")
        else:
            draw.text((start_x, txt['y']), content, fill=txt['color'], font=font, anchor="ls")        
    return img

# ==========================================
# 🎨 UI - STREAMLIT APP
# ==========================================
st.set_page_config(page_title="Auto Cert Pro", layout="wide")

# ตั้งค่า Session State
if "click_x" not in st.session_state: st.session_state.click_x = 0
if "click_y" not in st.session_state: st.session_state.click_y = 0
if 'texts' not in st.session_state: st.session_state.texts = []
if 'fonts_dict' not in st.session_state: st.session_state.fonts_dict = {}
if 'font_names' not in st.session_state: st.session_state.font_names = []

st.title("📜 Certificate Generator")

# --- SIDEBAR ---
with st.sidebar:
    st.header("1️⃣ อัปโหลดไฟล์")
    
    # 1. Template
    template_file = st.file_uploader("1. พื้นหลังเกียรติบัตร (JPG/PNG)", type=['jpg', 'jpeg', 'png'])
    if template_file:
        st.session_state.template = Image.open(template_file)

    # 2. Font
    font_file = st.file_uploader("2. ฟอนต์ภาษาไทย (.ttf) อัปโหลดเพิ่มได้", type=['ttf'])
    if font_file:
        st.session_state.font_bytes = font_file.getvalue() 
        f_name = font_file.name.split('.')[0]
        if f_name not in st.session_state.fonts_dict:
            st.session_state.fonts_dict[f_name] = font_file.getvalue()
            st.session_state.font_names.append(f_name)
            st.success(f"✅ โหลดฟอนต์ '{f_name}' สำเร็จ")
    else:
        if not st.session_state.fonts_dict:
            st.warning("⚠️ แนะนำให้อัปโหลดไฟล์ .ttf เพื่อให้ปรับขนาด/แสดงภาษาไทยได้")

    font_version = st.radio("ชนิดฟอนต์ (แก้ปัญหาสระลอย/กลายเป็นสี่เหลี่ยม)", 
                            ["ใหม่", "เก่า"], 
                            format_func=lambda x: "ฟอนต์รุ่นใหม่ (มาตรฐาน OpenType)" if x == "ใหม่" else "ฟอนต์รุ่นเก่า (ระบบ PUA ดั้งเดิม)")
    st.session_state.font_version = font_version

    # 3. Data
    data_file = st.file_uploader("3. รายชื่อ (ไฟล์ Excel/CSV)", type=['xlsx', 'xls', 'csv'])
    if data_file:
        if data_file.name.endswith('.csv'):
            st.session_state.data = pd.read_csv(data_file)
        else:
            st.session_state.data = pd.read_excel(data_file)

if 'template' not in st.session_state:
    st.info("👈 กรุณาอัปโหลด 'ต้นแบบเกียรติบัตร' ที่เมนูด้านซ้ายเพื่อเริ่มต้น")
    st.stop()

# --- MAIN AREA ---
st.header("2️⃣ กำหนดตำแหน่ง")

col_img, col_form = st.columns([1.5, 1])

with col_img:
    st.markdown("**🖱️ คลิกลงบนรูปภาพเพื่อดึงพิกัด(คลิกบนรูป หรือ ระบุ X และ Y)**")
    
    original_w, original_h = st.session_state.template.size
    display_w = 700 
    
    if original_w > display_w:
        ratio = original_w / display_w
        display_img = st.session_state.template.resize((display_w, int(original_h / ratio)))
    else:
        ratio = 1.0
        display_img = st.session_state.template

    coords = streamlit_image_coordinates(display_img, key="target_clicker")
    
    if coords is not None:
        st.session_state.click_x = int(coords['x'] * ratio)
        st.session_state.click_y = int(coords['y'] * ratio)

with col_form:
    st.markdown("**📝 ตั้งค่าข้อความ**")
    with st.form("add_text_form", clear_on_submit=False):
        t_type = st.radio("ชนิดข้อมูล", ["พิมพ์เอง", "ดึงจาก Excel"], horizontal=True)
        
        t_val, t_col = "", ""
        if "พิมพ์เอง" in t_type:
            t_val = st.text_input("ระบุข้อความ")
        else:
            if 'data' in st.session_state:
                t_col = st.selectbox("เลือกหัวข้อ (Column)", st.session_state.data.columns)
            else:
                st.warning("อัปโหลด Excel ก่อนครับ")

        c1, c2 = st.columns(2)
        x_pos = c1.number_input("แกน X", value=st.session_state.click_x)
        y_pos = c2.number_input("แกน Y", value=st.session_state.click_y)
        
        f_size = st.slider("ขนาดฟอนต์", 10, 500, value=60)
        f_color = st.color_picker("เลือกสี", value="#000000")
        
        selected_font = None
        if st.session_state.font_names:
            selected_font = st.selectbox("เลือกฟอนต์สำหรับข้อความนี้", st.session_state.font_names)
        
        if st.form_submit_button("➕ แทรกข้อความลงเกียรติบัตร"):
            if not st.session_state.fonts_dict:
                st.warning("อัปโหลดฟอนต์เพื่อปรับขนาด(หากไม่ทำจะปรับขนาดไม่ได้ แต่ระบบจะพยายามดึงฟอนต์ระบบมาใช้แทน)")
            
            st.session_state.texts.append({
                'type': 'static' if "พิมพ์เอง" in t_type else 'excel',
                'text': t_val, 'column': t_col,
                'x': x_pos, 'y': y_pos,
                'size': f_size, 'color': f_color,
                'font_name': selected_font 
            })
            st.rerun()

st.markdown("---")

# --- พรีวิวและจัดการข้อความ ---
st.header("3️⃣ ดูตัวอย่าง (Preview)")

if st.session_state.texts:
    for i, t in enumerate(st.session_state.texts):
        lbl = t['text'] if t['type'] == 'static' else f"จาก: {t['column']}"
        f_lbl = f" | ฟอนต์: {t.get('font_name', 'Default')}" if st.session_state.font_names else ""
        cols = st.columns([4, 1])
        cols[0].write(f"**{i+1}. {lbl}**{f_lbl} | ขนาด: {t['size']} | พิกัด: ({t['x']}, {t['y']})")
        if cols[1].button("🗑️ ลบ", key=f"del_{i}"):
            st.session_state.texts.pop(i)
            st.rerun()

    preview_row = None
    if 'data' in st.session_state:
        row_idx = st.number_input("ดูตัวอย่างแถวที่:", 0, max(0, len(st.session_state.data)-1), 0)
        preview_row = st.session_state.data.iloc[row_idx].to_dict()
    
    preview_img = render_certificate(st.session_state.template, st.session_state.texts, preview_row, st.session_state.font_version)
    st.image(preview_img, width=650)
else:
    st.info("ตั้งค่าข้อความด้านบนก่อนครับ")

# --- Export ---
if 'data' in st.session_state and st.session_state.texts:
    st.markdown("---")
    st.header("4️⃣ สร้างไฟล์ทั้งหมด")
    
    col_export1, col_export2 = st.columns(2)
    filename_col = col_export1.selectbox("เลือกคอลัมน์ชื่อไฟล์ (สำหรับ PNG/JPG/PDF)", st.session_state.data.columns)
    export_format = col_export2.radio("เลือกนามสกุลไฟล์", ["PNG", "JPG", "PDF", "PowerPoint"], horizontal=True)
    
    if export_format == "PowerPoint":
        st.info("💡 หมายเหตุ: ไฟล์ PowerPoint จะแยกกล่องข้อความให้คลิกแก้ไขได้ (ฟอนต์จะแสดงผลถูกต้องหากเครื่องคอมพิวเตอร์ที่เปิดไฟล์ติดตั้งฟอนต์นั้นไว้แล้ว)")

    if st.button("🚀 สร้างไฟล์", type="primary"):
        
        # ==========================================
        # 1. ส่งออกแบบ PowerPoint (ข้อความแก้ไขได้)
        # ==========================================
        if export_format == "PowerPoint":
            with st.spinner("กำลังสร้างไฟล์ PowerPoint (แยกข้อความ)..."):
                prs = Presentation()
                
                # กำหนดขนาดสไลด์ให้พอดีกับภาพต้นฉบับ (1 pixel = 9525 EMUs)
                img_width, img_height = st.session_state.template.size
                prs.slide_width = img_width * 9525
                prs.slide_height = img_height * 9525
                
                # ใช้ layout สไลด์เปล่า (ไม่มีกล่องข้อความขยะ)
                blank_slide_layout = prs.slide_layouts[6]
                
                # บันทึก "ภาพพื้นหลังเปล่าๆ" ลงหน่วยความจำ 
                # (ป้องกันข้อผิดพลาดกรณีเป็นไฟล์ที่มีความโปร่งใส ต้องแปลงเป็น RGB ก่อน)
                bg_io = BytesIO()
                bg_img = st.session_state.template.convert('RGB') if st.session_state.template.mode != 'RGB' else st.session_state.template
                bg_img.save(bg_io, format="PNG")
                
                for idx, row in st.session_state.data.iterrows():
                    slide = prs.slides.add_slide(blank_slide_layout)
                    
                    # 1. แปะภาพพื้นหลัง
                    bg_io.seek(0)
                    slide.shapes.add_picture(bg_io, 0, 0, width=prs.slide_width, height=prs.slide_height)
                    
                    # 2. สร้าง Text Box วางซ้อนทับตามพิกัดของข้อความแต่ละตัว
                    for txt in st.session_state.texts:
                        # ดึงเนื้อหา
                        if txt['type'] == 'static':
                            content = txt['text']
                        else:
                            if txt['column'] in row:
                                val = row[txt['column']]
                                content = str(val) if pd.notna(val) else ""
                            else:
                                continue
                                
                        if not content: continue
                        
                        # คำนวณขนาดและตำแหน่ง (แปลงจาก px เป็น pt และ EMUs)
                        # ขนาดฟอนต์ PowerPoint ใช้พอยต์ (1 px = 0.75 pt)
                        font_pt_size = txt['size'] * 0.75 
                        
                        # จุด Y ในโปรแกรมเราคือฐานบรรทัด แต่ Text Box ใน PPTX คือจุดบนสุด 
                        # จึงต้องเลื่อน Y ขึ้นไปตามขนาดฟอนต์
                        pixel_y = txt['y'] - txt['size']
                        
                        # จัดให้อยู่กึ่งกลางเป๊ะๆ: สร้าง Text Box ให้กว้างสุด แล้วตั้ง Align เป็น Center
                        box_width_px = img_width 
                        left_px = txt['x'] - (box_width_px / 2)
                        
                        left = int(left_px * 9525)
                        top = int(pixel_y * 9525)
                        width = int(box_width_px * 9525)
                        height = int(txt['size'] * 1.5 * 9525) # เผื่อความสูงให้ไม่ตกขอบ
                        
                        txBox = slide.shapes.add_textbox(left, top, width, height)
                        tf = txBox.text_frame
                        tf.word_wrap = False
                        
                        # ตั้งค่าพารากราฟและใส่ข้อความ
                        p = tf.paragraphs[0]
                        p.alignment = PP_ALIGN.CENTER
                        p.text = content
                        
                        # กำหนดรูปแบบฟอนต์
                        run = p.runs[0]
                        run.font.size = Pt(font_pt_size)
                        
                        # แปลงสี Hex เป็น RGBColor สำหรับ PowerPoint
                        hex_color = txt['color'].lstrip('#')
                        if len(hex_color) == 6:
                            run.font.color.rgb = RGBColor.from_string(hex_color)
                            
                        # ตั้งชื่อฟอนต์
                        if txt.get('font_name'):
                            run.font.name = txt['font_name']
                
                # เตรียมดาวน์โหลด
                pptx_io = BytesIO()
                prs.save(pptx_io)
                pptx_io.seek(0)
                
                st.success("✅ สร้างไฟล์ PowerPoint เรียบร้อย (คลิกแก้ไขข้อความข้างในได้เลย)!")
                st.download_button(
                    label="📥 ดาวน์โหลดไฟล์ PowerPoint (.pptx)", 
                    data=pptx_io, 
                    file_name="certificates.pptx", 
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

        # ==========================================
        # 2. ส่งออกแบบปกติ (ไฟล์ ZIP ประกอบด้วย PNG, JPG หรือ PDF)
        # ==========================================
        else:
            zip_buffer = BytesIO()
            with st.spinner(f"กำลังสร้างไฟล์เกียรติบัตร ({export_format})..."):
                with zipfile.ZipFile(zip_buffer, 'w') as zf:
                    for idx, row in st.session_state.data.iterrows():
                        final_img = render_certificate(st.session_state.template, st.session_state.texts, row.to_dict(), st.session_state.font_version)
                        img_io = BytesIO()
                        
                        pil_format = "JPEG" if export_format == "JPG" else export_format
                        final_img.save(img_io, format=pil_format, resolution=300.0)
                        
                        clean_name = sanitize_filename(row[filename_col])
                        file_extension = export_format.lower()
                        
                        zf.writestr(f"{clean_name}.{file_extension}", img_io.getvalue())
                
                st.success("✅ สร้างไฟล์ทั้งหมดสำเร็จ!")
                st.download_button(f"📥 ดาวน์โหลดไฟล์ ZIP ({export_format})", zip_buffer.getvalue(), "certificates.zip", "application/zip")
